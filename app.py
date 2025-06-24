import streamlit as st
import os
import tempfile
import requests
import zipfile
import io
import shutil
import nltk
import pandas as pd
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.chains import ConversationalRetrievalChain
from langchain.document_loaders import (
    PyPDFLoader,
    TextLoader,
    Docx2txtLoader
)
from langchain.schema import Document
from langchain.memory import ConversationBufferMemory
import glob
import json
from typing import List

# Set NLTK_DATA path to a writable directory
try:
    nltk_data_dir = os.path.join(tempfile.gettempdir(), "nltk_data")
    os.makedirs(nltk_data_dir, exist_ok=True)
    nltk.data.path.append(nltk_data_dir)
    os.environ["NLTK_DATA"] = nltk_data_dir
    # Download necessary NLTK data without any stdout output
    nltk.download('punkt', download_dir=nltk_data_dir, quiet=True)
    nltk.download('averaged_perceptron_tagger', download_dir=nltk_data_dir, quiet=True)
except Exception:
    pass  # Silently handle NLTK setup issues

# Page configuration
st.set_page_config(
    page_title="Kaizen Engineers Knowledge Base",
    page_icon="📚",
    layout="wide"
)

# GitHub ZIP URL (hidden from frontend)
GITHUB_ZIP_URL = "https://github.com/pravinraut05/Dataset/raw/main/orders.zip"

# Initialize session state variables
if "conversation" not in st.session_state:
    st.session_state.conversation = None
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "processed_data" not in st.session_state:
    st.session_state.processed_data = False
if "temp_dir" not in st.session_state:
    st.session_state.temp_dir = None
if "initialization_started" not in st.session_state:
    st.session_state.initialization_started = False

# Function to download and extract ZIP file from GitHub
def download_and_extract_zip(github_zip_url):
    try:
        # Create temporary directory
        temp_dir = tempfile.mkdtemp()
        
        # Download the ZIP file (silently)
        response = requests.get(github_zip_url)
        if response.status_code != 200:
            return None
            
        # Extract the ZIP file (silently)
        with zipfile.ZipFile(io.BytesIO(response.content)) as zip_ref:
            zip_ref.extractall(temp_dir)
        
        return temp_dir
    except Exception:
        return None

# Enhanced CSV processing function
def process_csv_file(csv_path: str) -> List[Document]:
    """
    Process CSV file to create structured documents with better context preservation
    """
    documents = []
    
    try:
        # Read CSV file
        df = pd.read_csv(csv_path, encoding='utf-8')
    except UnicodeDecodeError:
        try:
            df = pd.read_csv(csv_path, encoding='latin-1')
        except Exception:
            return documents
    except Exception:
        return documents
    
    if df.empty:
        return documents
    
    # Get basic info about the dataset
    filename = os.path.basename(csv_path)
    num_rows, num_cols = df.shape
    columns = list(df.columns)
    
    # Create a summary document
    summary_content = f"""
    Dataset: {filename}
    Number of rows: {num_rows}
    Number of columns: {num_cols}
    Columns: {', '.join(columns)}
    
    Dataset Overview:
    {df.describe(include='all').to_string() if not df.empty else 'No data to describe'}
    
    Sample Data (first 5 rows):
    {df.head().to_string(index=False) if not df.empty else 'No data available'}
    """
    
    summary_doc = Document(
        page_content=summary_content,
        metadata={"source": csv_path, "type": "csv_summary", "filename": filename}
    )
    documents.append(summary_doc)
    
    # Create chunks for better retrieval - process in batches
    chunk_size = 50  # Number of rows per chunk
    
    for i in range(0, len(df), chunk_size):
        chunk_df = df.iloc[i:i+chunk_size]
        
        # Create structured content for this chunk
        chunk_content = f"""
        Dataset: {filename} (Rows {i+1}-{min(i+chunk_size, len(df))})
        
        Data:
        {chunk_df.to_string(index=False)}
        
        Row-by-row details:
        """
        
        # Add detailed row information
        for idx, row in chunk_df.iterrows():
            row_details = []
            for col in df.columns:
                if pd.notna(row[col]):
                    row_details.append(f"{col}: {row[col]}")
            
            if row_details:
                chunk_content += f"\nRow {idx + 1}: {', '.join(row_details)}"
        
        # Add column-wise summaries for this chunk
        chunk_content += "\n\nColumn summaries for this chunk:"
        for col in df.columns:
            col_data = chunk_df[col].dropna()
            if not col_data.empty:
                if pd.api.types.is_numeric_dtype(col_data):
                    chunk_content += f"\n{col}: min={col_data.min()}, max={col_data.max()}, mean={col_data.mean():.2f}"
                else:
                    unique_vals = col_data.unique()[:5]  # Show first 5 unique values
                    chunk_content += f"\n{col}: sample values: {', '.join(map(str, unique_vals))}"
        
        chunk_doc = Document(
            page_content=chunk_content,
            metadata={
                "source": csv_path, 
                "type": "csv_chunk", 
                "filename": filename,
                "chunk_start": i,
                "chunk_end": min(i+chunk_size, len(df))
            }
        )
        documents.append(chunk_doc)
    
    # Create documents for each column with its complete data
    for col in df.columns:
        col_data = df[col].dropna()
        if not col_data.empty:
            col_content = f"""
            Column: {col} from dataset {filename}
            Data type: {str(df[col].dtype)}
            Non-null values: {len(col_data)} out of {len(df)}
            
            """
            
            if pd.api.types.is_numeric_dtype(col_data):
                col_content += f"""
                Statistical summary:
                - Minimum: {col_data.min()}
                - Maximum: {col_data.max()}
                - Mean: {col_data.mean():.2f}
                - Median: {col_data.median()}
                - Standard deviation: {col_data.std():.2f}
                
                All values: {', '.join(map(str, col_data.tolist()))}
                """
            else:
                unique_values = col_data.unique()
                col_content += f"""
                Unique values count: {len(unique_values)}
                Unique values: {', '.join(map(str, unique_values))}
                
                All values: {', '.join(map(str, col_data.tolist()))}
                """
            
            col_doc = Document(
                page_content=col_content,
                metadata={
                    "source": csv_path, 
                    "type": "csv_column", 
                    "filename": filename,
                    "column_name": col
                }
            )
            documents.append(col_doc)
    
    return documents

# Enhanced Excel processing function
def process_excel_file(excel_path: str) -> List[Document]:
    """
    Process Excel file similar to CSV with support for multiple sheets
    """
    documents = []
    
    try:
        # Load all sheets from the Excel file
        xl = pd.ExcelFile(excel_path)
        filename = os.path.basename(excel_path)
        
        # Create a summary document for the entire Excel file
        summary_content = f"""
        Excel File: {filename}
        Number of sheets: {len(xl.sheet_names)}
        Sheet names: {', '.join(xl.sheet_names)}
        """
        
        summary_doc = Document(
            page_content=summary_content,
            metadata={"source": excel_path, "type": "excel_summary", "filename": filename}
        )
        documents.append(summary_doc)
        
        # Process each sheet
        for sheet_name in xl.sheet_names:
            try:
                df = pd.read_excel(excel_path, sheet_name=sheet_name)
                if df.empty:
                    continue
                
                # Process this sheet similar to CSV
                sheet_documents = process_dataframe_as_csv(df, f"{excel_path}#{sheet_name}", f"{filename}#{sheet_name}")
                documents.extend(sheet_documents)
                
            except Exception:
                continue
                
    except Exception:
        pass
    
    return documents

def process_dataframe_as_csv(df: pd.DataFrame, source_path: str, display_name: str) -> List[Document]:
    """
    Helper function to process a DataFrame similar to CSV processing
    """
    documents = []
    
    if df.empty:
        return documents
    
    num_rows, num_cols = df.shape
    columns = list(df.columns)
    
    # Create a summary document
    summary_content = f"""
    Dataset: {display_name}
    Number of rows: {num_rows}
    Number of columns: {num_cols}
    Columns: {', '.join(columns)}
    
    Dataset Overview:
    {df.describe(include='all').to_string() if not df.empty else 'No data to describe'}
    
    Sample Data (first 5 rows):
    {df.head().to_string(index=False) if not df.empty else 'No data available'}
    """
    
    summary_doc = Document(
        page_content=summary_content,
        metadata={"source": source_path, "type": "data_summary", "filename": display_name}
    )
    documents.append(summary_doc)
    
    # Create chunks for better retrieval
    chunk_size = 50
    
    for i in range(0, len(df), chunk_size):
        chunk_df = df.iloc[i:i+chunk_size]
        
        chunk_content = f"""
        Dataset: {display_name} (Rows {i+1}-{min(i+chunk_size, len(df))})
        
        Data:
        {chunk_df.to_string(index=False)}
        
        Row-by-row details:
        """
        
        for idx, row in chunk_df.iterrows():
            row_details = []
            for col in df.columns:
                if pd.notna(row[col]):
                    row_details.append(f"{col}: {row[col]}")
            
            if row_details:
                chunk_content += f"\nRow {idx + 1}: {', '.join(row_details)}"
        
        chunk_doc = Document(
            page_content=chunk_content,
            metadata={
                "source": source_path, 
                "type": "data_chunk", 
                "filename": display_name,
                "chunk_start": i,
                "chunk_end": min(i+chunk_size, len(df))
            }
        )
        documents.append(chunk_doc)
    
    return documents

# Function to load documents from directory with enhanced CSV handling
def load_documents(directory_path):
    documents = []
    
    # Handle PDF files
    for pdf_path in glob.glob(f"{directory_path}/**/*.pdf", recursive=True):
        try:
            loader = PyPDFLoader(pdf_path)
            docs = loader.load()
            documents.extend(docs)
        except Exception:
            pass
    
    # Handle text files
    for text_path in glob.glob(f"{directory_path}/**/*.txt", recursive=True):
        try:
            loader = TextLoader(text_path, encoding='utf-8')
            docs = loader.load()
            documents.extend(docs)
        except UnicodeDecodeError:
            try:
                loader = TextLoader(text_path, encoding='latin-1')
                docs = loader.load()
                documents.extend(docs)
            except Exception:
                pass
        except Exception:
            pass
    
    # Handle Markdown files
    for md_path in glob.glob(f"{directory_path}/**/*.md", recursive=True):
        try:
            loader = TextLoader(md_path, encoding='utf-8')
            docs = loader.load()
            documents.extend(docs)
        except Exception:
            pass
    
    # Handle DOCX files
    for docx_path in glob.glob(f"{directory_path}/**/*.docx", recursive=True):
        try:
            loader = Docx2txtLoader(docx_path)
            docs = loader.load()
            documents.extend(docs)
        except Exception:
            pass
    
    # Handle DOC files
    for doc_path in glob.glob(f"{directory_path}/**/*.doc", recursive=True):
        try:
            with open(doc_path, 'rb') as f:
                content = f.read().decode('latin-1', errors='replace')
                doc = Document(page_content=content, metadata={"source": doc_path})
                documents.append(doc)
        except Exception:
            pass
    
    # Enhanced CSV handling
    for csv_path in glob.glob(f"{directory_path}/**/*.csv", recursive=True):
        try:
            csv_documents = process_csv_file(csv_path)
            documents.extend(csv_documents)
        except Exception:
            pass
    
    # Enhanced Excel handling
    for excel_path in glob.glob(f"{directory_path}/**/*.xlsx", recursive=True) + glob.glob(f"{directory_path}/**/*.xls", recursive=True):
        try:
            excel_documents = process_excel_file(excel_path)
            documents.extend(excel_documents)
        except Exception:
            pass
    
    # Handle PowerPoint files
    for ppt_path in glob.glob(f"{directory_path}/**/*.pptx", recursive=True) + glob.glob(f"{directory_path}/**/*.ppt", recursive=True):
        try:
            from pptx import Presentation
            
            try:
                prs = Presentation(ppt_path)
                content = []
                
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    
                    if slide_text:
                        content.append(f"Slide {i+1}: {' '.join(slide_text)}")
                
                if content:
                    full_content = "\n\n".join(content)
                    doc = Document(page_content=full_content, metadata={"source": ppt_path})
                    documents.append(doc)
            except:
                doc = Document(
                    page_content=f"PowerPoint file found at {ppt_path}. This is a binary file that couldn't be fully processed.",
                    metadata={"source": ppt_path}
                )
                documents.append(doc)
        except Exception:
            pass
    
    return documents

# Function to split documents into chunks with improved parameters for structured data
def split_documents(documents):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=2000,  # Increased chunk size for better context
        chunk_overlap=400,  # Increased overlap for better continuity
        length_function=len,
        separators=["\n\n", "\n", " ", ""]  # Better separators for structured data
    )
    chunks = text_splitter.split_documents(documents)
    return chunks

# Function to create vector database from chunks
def create_vectordb(chunks, openai_api_key):
    try:
        os.environ["OPENAI_API_KEY"] = openai_api_key
        
        embeddings = OpenAIEmbeddings(
            model="text-embedding-ada-002",
            openai_api_key=openai_api_key,
        )
        
        if not chunks:
            return None
            
        vectorstore = FAISS.from_documents(chunks, embeddings)
        return vectorstore
    except Exception:
        return None

# Function to create conversation chain with enhanced configuration
def create_conversation_chain(vectorstore):
    llm = ChatOpenAI(
        temperature=0.1,  # Slightly increased for more natural responses
        model_name="gpt-3.5-turbo-16k",  # Use 16k model for better context
        openai_api_key=os.environ.get("OPENAI_API_KEY")
    )
    
    memory = ConversationBufferMemory(
        memory_key="chat_history", 
        return_messages=True,
        output_key="answer"
    )
    
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(
            search_type="mmr",  # Use MMR for better diversity
            search_kwargs={
                "k": 8,  # Increased number of documents to retrieve
                "fetch_k": 20,  # Fetch more candidates for MMR
                "lambda_mult": 0.7  # Balance between relevance and diversity
            }
        ),
        memory=memory,
        verbose=False,
        return_source_documents=True,
        combine_docs_chain_kwargs={
            "prompt": None  # Use default prompt which works well with structured data
        }
    )
    
    return conversation_chain

# Function to handle user input with enhanced response formatting
def handle_user_input(user_question):
    with st.chat_message("user"):
        st.write(user_question)
    
    with st.chat_message("assistant"):
        with st.spinner("Thinking..."):
            response = st.session_state.conversation({"question": user_question})
            st.session_state.chat_history = response['chat_history']
            
            # Display the main answer
            st.write(response['answer'])
            
            # Optionally show source information for debugging
            if 'source_documents' in response and response['source_documents']:
                with st.expander("📄 Sources", expanded=False):
                    for i, doc in enumerate(response['source_documents'][:3]):  # Show top 3 sources
                        source = doc.metadata.get('source', 'Unknown')
                        filename = os.path.basename(source)
                        st.write(f"**Source {i+1}:** {filename}")
                        if 'type' in doc.metadata:
                            st.write(f"**Type:** {doc.metadata['type']}")

# Function to initialize the knowledge base
def initialize_knowledge_base(openai_api_key):
    temp_dir = download_and_extract_zip(GITHUB_ZIP_URL)
    
    if temp_dir:
        st.session_state.temp_dir = temp_dir
        
        try:
            documents = load_documents(temp_dir)
            
            if documents and len(documents) > 0:
                chunks = split_documents(documents)
                vectorstore = create_vectordb(chunks, openai_api_key)
                
                if vectorstore:
                    st.session_state.conversation = create_conversation_chain(vectorstore)                    
                    st.session_state.processed_data = True
            else:
                if st.session_state.temp_dir and os.path.exists(st.session_state.temp_dir):
                    shutil.rmtree(st.session_state.temp_dir)
                st.session_state.temp_dir = None
        except Exception:
            if st.session_state.temp_dir and os.path.exists(st.session_state.temp_dir):
                shutil.rmtree(st.session_state.temp_dir)
            st.session_state.temp_dir = None

# Main app title and description
st.title("📚 Kaizen Engineers Knowledge Base")
st.markdown("""
This application provides access to Kaizen Engineers' knowledge resources using AI-powered search and retrieval.
Enhanced with improved CSV and structured data processing for better accuracy.
""")

# Sidebar for company information
with st.sidebar:
    # Company logo handling
    logo_found = False
    possible_logo_paths = [
        "logo.png",
        os.path.join(os.getcwd(), "logo.png"),
        os.path.abspath("logo.png"),
        "/app/logo.png",
        os.path.join(os.path.dirname(__file__), "logo.png"),
    ]
    
    for logo_path in possible_logo_paths:
        try:
            if os.path.exists(logo_path) and os.path.isfile(logo_path):
                st.image(logo_path, use_container_width=True)
                logo_found = True
                break
        except Exception:
            continue
    
    if not logo_found:
        try:
            st.markdown("""
            <style>
            .kaizen-logo {
                display: block;
                margin-left: auto;
                margin-right: auto;
                width: 100%;
                max-width: 200px;
                height: auto;
                text-align: center;
                font-family: Arial, sans-serif;
                color: #2C3E50;
            }
            .kaizen-logo-image {
                background: linear-gradient(135deg, #2980B9, #27AE60);
                width: 80px;
                height: 80px;
                border-radius: 10px;
                margin: 0 auto 10px auto;
                position: relative;
                display: flex;
                align-items: center;
                justify-content: center;
                color: white;
                font-size: 28px;
                font-weight: bold;
            }
            </style>
            <div class="kaizen-logo">
                <div class="kaizen-logo-image">K</div>
                <div style="text-align: center; font-weight: bold;">Kaizen Engineers</div>
            </div>
            """, unsafe_allow_html=True)
        except:
            pass
    
    st.markdown("Created by Udayam AI Labs 🚀")
    
    if st.session_state.processed_data and st.session_state.conversation:
        if st.button("Clear Chat", key="clear_chat_button", help="Clear current chat history"):
            st.session_state.chat_history = []
            if st.session_state.conversation:
                st.session_state.conversation.memory.clear()

# API key handling and initialization
openai_api_key = os.environ.get("OPENAI_API_KEY", "")
if not openai_api_key:
    st.warning("Please enter your OpenAI API key to continue")
    openai_api_key = st.text_input("OpenAI API Key", type="password")
    if openai_api_key:
        os.environ["OPENAI_API_KEY"] = openai_api_key

if openai_api_key and not st.session_state.initialization_started:
    st.session_state.initialization_started = True
    with st.spinner("Loading knowledge base..."):
        initialize_knowledge_base(openai_api_key)

# Main chat interface
if st.session_state.processed_data:
    st.header("Chat with Kaizen Engineers Knowledge Base")
    
    # Display chat history
    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            with st.chat_message("user"):
                st.write(message.content)
        else:
            with st.chat_message("assistant"):
                st.write(message.content)
    
    # User input
    user_question = st.chat_input("Ask a question about Kaizen Engineers...")
    if user_question:
        handle_user_input(user_question)
else:
    if st.session_state.initialization_started and openai_api_key:
        st.info("Setting up the knowledge base with enhanced CSV processing. This may take a moment...")
    elif not openai_api_key:
        st.header("Sample Questions You Can Ask")
        st.markdown("""
        Once the knowledge base is initialized, you can ask questions like:
        - What are Kaizen Engineers' main areas of expertise?
        - What projects has Kaizen Engineers completed recently?
        - What services does Kaizen Engineers offer?
        - Who are the key personnel at Kaizen Engineers?
        - What is Kaizen Engineers' approach to quality control?
        - **For CSV files:** Show me the summary of the data, What are the column names?, What's the average of [column name]?, How many records are there?
        """)
